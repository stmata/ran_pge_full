import tempfile
import re
from moviepy.editor import VideoFileClip
import os
from PyPDF2 import PdfReader
from pptx import Presentation


class Extractor:
    def __init__(self):
        pass

    def extract_audio(self, video_file, audio_path: str) -> None:
        """
        Extracts audio from a video file and saves it to the specified audio path.

        :param video_file: An object containing the video file (Flask FileStorage object).
        :param audio_path: The path where the extracted audio file will be saved.
        """
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as temp_video_file:
                video_file.save(temp_video_file.name)
                print(temp_video_file.name)
                with VideoFileClip(temp_video_file.name) as video:
                    audio = video.audio
                    print(audio)
                    audio.write_audiofile(audio_path)

        except Exception as e:
            print(f"Une erreur est survenue lors de l'extraction audio : {e}")
        finally:
            if 'audio' in locals():
                audio.close()
            if 'video' in locals():
                video.close()
            os.remove(temp_video_file.name)


    def extract_youtube_id(self, url: str) -> str:
        """
        Extracts YouTube video ID from the given URL.

        :param url: URL of the YouTube video.
        :return: YouTube video ID or None if not found.
        """
        regex = r"(youtu\.be\/|v=)([^&]+)"
        match = re.search(regex, url)
        if match:
            return match.group(2)
        return None
    
    def pdf_to_text(self, pdf_file):
        try:
            content = ""
            with open(pdf_file, "rb") as file:
                reader = PdfReader(file)
                for page in reader.pages:
                    text = page.extract_text() if page.extract_text() else ""
                    text = re.sub(r'[^a-zA-Z0-9\s]', '', text) 
                    text = re.sub(r'\s+', ' ', text)  
                    text = text.replace('\n', ' ')  
                    content += text
            return content
        except Exception as e:
            return f"Error extracting text from PDF: {str(e)}"
        
    def extract_text_from_pptx(self, pptx_file):
        try:
            presentation = Presentation(pptx_file)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text = shape.text
                        slide_text = re.sub(r'[^a-zA-Z0-9\s]', '', slide_text)  
                        slide_text = re.sub(r'\s+', ' ', slide_text)  
                        text.append(slide_text)
            return " ".join(text).strip()  
        except Exception as e:
            print(f"An error occurred while extracting text from PPTX: {e}")
            return None
